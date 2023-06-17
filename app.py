import streamlit as st
import sys
import PyPDF2
from pptx import Presentation

sys.path.insert(1, 'Eslam')
sys.path.insert(2, 'mariam')
from Eslam_main import GenerateSummary

st.set_page_config(page_title="Study Sync", page_icon="https://github.com/AnasMations/StudySync/blob/36ae1cad78544b8a07239d1cefa141a59f6305c8/img/icon.png?raw=true")

def style():
    st.markdown(
        """
        <style>
            
            [data-testid="stAppViewContainer"] {
                background: linear-gradient(220deg, #fbeaf7,#ffffff, #dfeefb);
                color: black;
            }
            

        div.stButton > button:first-child {
            font-size: 24px;
            padding: 16px 32px;
            background-color: #FFFFFF;
            color: #262730;
            border: none;
            display: block;
            margin-left: auto;
            margin-right:auto;
            border-radius: 12px; /* Add rounded edges */
            box-shadow: 0 0px 15px rgba(0, 0, 0, 0.2); /* Add shadow for a 3D look */
            transition: all 0.3s ease-in-out; /* Add smooth transition for hover effects */
        
        }

        div.stButton > button:first-child:hover {
            cursor: pointer;
            border: none;
            background-color: #6ebabf;
            transform: translateY(-5px) scale(1.02); /* Animate movement upwards and scale */
            box-shadow: 0 10px 40px rgba(110, 186, 191, 1); /* Enlarge shadow on hover for a 3D effect */
        
        }
        
        @import url('https://fonts.googleapis.com/css2?family=Poppins&display=swap');

        * {
            font-family: 'Poppins', sans-serif;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

############################
# SECTION DEFINATIONS
############################
def HEADER():
    '''
    This function displays the header section of the web app.
    '''

    # Create a container layout
    container = st.container()

    # Add some CSS styles using markdown to position and center the item horizontally
    container.markdown(
        """
        <style>
        .logo {{
            position: absolute;
            top: 0;
            left: 50%;
            margin-left: -{width}px;
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )
    with container:
        st.image("img/Study Sync.png")

    st.subheader("Your one-stop solution for studying smarter, not harder.")

def FEATURES():
    '''
    This function displays the features section of the web app.
    '''

    # Style for the individual feature boxes
    style = """
    <style>
        .feature-box {{
            background-color: #fafafa;
            border: 1px solid #ccc;
            border-radius: 5px;
            margin-bottom: 10px;
            padding: 15px;
        }}
    </style>
    """

    st.subheader("Features")
    st.markdown(style, unsafe_allow_html=True)

    features = [
        {"title": "User-friendly Interface", "emoji": "🌟", "desc": "Study Sync offers an intuitive and engaging user interface designed using Streamlit. It ensures a smooth and enjoyable user experience throughout the application."},
        {"title": "Study Material Input", "emoji": "📚", "desc": "Users can easily input their study materials through various options, including text input or file upload. This flexibility allows users to seamlessly integrate their existing study resources into Study Sync."},
        {"title": "Automatic Summarization", "emoji": "✂️", "desc": "Study Sync employs the power of ChatGPT to automatically generate concise summaries of study materials. This feature saves time and effort by condensing lengthy content into key points for quick comprehension."},
        {"title": "Flashcards Generation", "emoji": "🃏", "desc": "With Study Sync, users can create interactive flashcards from important concepts within their study materials. Flashcards serve as effective tools for reinforcing knowledge and facilitating active recall."},
        {"title": "Test Question Generation", "emoji": "✅", "desc": "Study Sync generates diverse test questions for self-assessment purposes. These questions cover various topics, ensuring comprehensive understanding and preparation for exams or assessments."},
        {"title": "Video Summaries", "emoji": "🎥", "desc": "Study Sync leverages the capabilities of Langchain to produce engaging video summaries of study materials. These video summaries provide an alternative medium for reviewing content and enhancing the learning experience."},
        {"title": "Mind Maps", "emoji": "🧠", "desc": "Study Sync leverages the capabilities of Langchain to produce Mind maps of study materials. These mind maps provide a structued outline for your study"}
    ]

    # Iterate through the features and display each one in a rectangle-box
    i = 0
    for feature in features:
        
        with st.expander(feature["emoji"]+" "+feature["title"]):
            st.markdown(f'<div class="feature-box"><h3>{feature["title"]}</h3><p>{feature["desc"]}</p></div>', unsafe_allow_html=True)



def GENERATE():
    """
    This function displays the get started section of the web app.
    """

    st.subheader("Get Started")

    # Dropdown for input_type
    input_options = ['Choose Input Type', 'Presentation slides', '.PDF', 'Video transcript', 'Question']
    input_selection = st.selectbox('Select Input Type:', input_options)
    input_type = input_options.index(input_selection)-1

    # Dropdown for output_type
    output_options = ['Choose Output Type', 'Questions and Answers', 'Multiple choice questions', 'Summarization', 'Bullet points']
    output_selection = st.selectbox('Select Output Type:', output_options)
    output_type = output_options.index(output_selection)-1

    if input_type == 0: #presentation         
        uploaded_file = st.file_uploader('')

        if uploaded_file is not None:

            content = ""
            presentation = Presentation(uploaded_file)

            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                content += run.text
            if st.button("Generate"):
                with st.spinner(text="Generating Summary..."):
                    st.write(GenerateSummary(input_type, output_type, content))
    
    elif input_type == 1: # Pdf
        
        uploaded_file = st.file_uploader('')

        if uploaded_file is not None:
            content = ""
            
            reader = PyPDF2.PdfReader(uploaded_file)
            num_pages = len(reader.pages)

            for page in range(num_pages):
                page_obj = reader.pages[page]
                content += page_obj.extract_text()
            if st.button("Generate"):
                with st.spinner(text="Generating Summary..."):
                    st.write(GenerateSummary(input_type, output_type, content))

    elif input_type == 2: # Video transcript
        url = st.text_input("Enter the URL of the video:")
        if url:
            if st.button("Generate"):
                with st.spinner(text="Generating Summary..."):
                    st.write(GenerateSummary(input_type, output_type, url))

    elif input_type == 3: # Question
        user_input = st.text_input("Enter study topic:")
        if user_input:
            if st.button("Generate"):
                with st.spinner(text="Generating Summary..."):
                    st.write(GenerateSummary(input_type, output_type, user_input))

############################
# CALL SECTIONS
############################
style()
HEADER()
GENERATE()
FEATURES()